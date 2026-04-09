#!/usr/bin/env python3
"""
Sankofa — African Regulatory Intelligence
Report Generator

Parsing backends (no API key required — all local):
  PDF   → PyMuPDF     (word-level bounding boxes + page screenshots)
  DOCX  → python-docx (paragraph text extraction)
  PPTX  → python-pptx (slide text extraction)
  Plain → read directly

Two modes:
  --parse-only : Discover and parse all supported files in a directory, output JSON
  --answer-file: Generate an HTML compliance report with visual citations
"""

import argparse
import base64
import html as html_module
import json
import os
import re
import sys
import time
import webbrowser
from concurrent.futures import ThreadPoolExecutor, as_completed
from datetime import datetime
from pathlib import Path

PDF_EXT  = {".pdf"}
DOCX_EXT = {".docx", ".docm"}
PPTX_EXT = {".pptx", ".pptm"}
PLAINTEXT_EXTENSIONS = {".txt", ".md", ".markdown", ".rst", ".text"}

JURISDICTION_NAMES = {
    "NG": "Nigeria",     "KE": "Kenya",        "GH": "Ghana",
    "ZA": "South Africa","EG": "Egypt",        "TZ": "Tanzania",
    "RW": "Rwanda",      "UG": "Uganda",       "SN": "Senegal",
    "ET": "Ethiopia",    "CI": "Côte d'Ivoire","CM": "Cameroon",
    "MZ": "Mozambique",  "ZM": "Zambia",       "MW": "Malawi",
    "ZW": "Zimbabwe",    "BW": "Botswana",     "NA": "Namibia",
    "MU": "Mauritius",   "MA": "Morocco",      "TN": "Tunisia",
    "DZ": "Algeria",
}
JURISDICTION_FLAGS = {
    "NG": "🇳🇬", "KE": "🇰🇪", "GH": "🇬🇭", "ZA": "🇿🇦",
    "EG": "🇪🇬", "TZ": "🇹🇿", "RW": "🇷🇼", "UG": "🇺🇬",
    "SN": "🇸🇳", "ET": "🇪🇹", "CI": "🇨🇮", "CM": "🇨🇲",
    "MZ": "🇲🇿", "ZM": "🇿🇲", "MW": "🇲🇼", "ZW": "🇿🇼",
    "BW": "🇧🇼", "NA": "🇳🇦", "MU": "🇲🇺", "MA": "🇲🇦",
    "TN": "🇹🇳", "DZ": "🇩🇿",
}


# ── File Discovery ────────────────────────────────────────────────────

def discover_files(data_dir: Path, max_files: int):
    pdf_files, docx_files, pptx_files, plaintext_files = [], [], [], []

    if not data_dir.is_dir():
        print(f"Error: {data_dir} is not a directory", file=sys.stderr)
        sys.exit(1)

    for f in sorted(data_dir.iterdir()):
        if not f.is_file():
            continue
        ext = f.suffix.lower()
        if ext in PDF_EXT:
            pdf_files.append(f)
        elif ext in DOCX_EXT:
            docx_files.append(f)
        elif ext in PPTX_EXT:
            pptx_files.append(f)
        elif ext in PLAINTEXT_EXTENSIONS:
            plaintext_files.append(f)

    all_files = pdf_files + docx_files + pptx_files + plaintext_files
    if len(all_files) > max_files:
        print(f"Warning: Found {len(all_files)} files, capping at {max_files}", file=sys.stderr)
        all_files = all_files[:max_files]
        pdf_files      = [f for f in all_files if f.suffix.lower() in PDF_EXT]
        docx_files     = [f for f in all_files if f.suffix.lower() in DOCX_EXT]
        pptx_files     = [f for f in all_files if f.suffix.lower() in PPTX_EXT]
        plaintext_files = [f for f in all_files if f.suffix.lower() in PLAINTEXT_EXTENSIONS]

    return pdf_files, docx_files, pptx_files, plaintext_files


# ── PDF parsing with PyMuPDF (word-level bounding boxes) ─────────────

def _parse_pdf(filepath: Path) -> dict:
    import fitz  # pymupdf

    t0 = time.perf_counter()
    pdf = fitz.open(str(filepath))
    pages_data = []

    for i in range(len(pdf)):
        page = pdf[i]
        # get_text("words") → list of (x0, y0, x1, y1, word, block_no, line_no, word_no)
        words = page.get_text("words")
        text_items = [
            {
                "text":   w[4],
                "x":      w[0],
                "y":      w[1],
                "width":  w[2] - w[0],
                "height": w[3] - w[1],
            }
            for w in words
        ]
        pages_data.append({
            "pageNum":   i + 1,
            "width":     page.rect.width,
            "height":    page.rect.height,
            "text":      page.get_text("text"),
            "textItems": text_items,
        })

    pdf.close()
    return {
        "name":      filepath.name,
        "path":      str(filepath),
        "type":      "liteparse",   # keep key name for HTML template compatibility
        "parseTime": round(time.perf_counter() - t0, 3),
        "pages":     pages_data,
    }


# ── DOCX parsing with python-docx ────────────────────────────────────

def _parse_docx(filepath: Path) -> dict:
    from docx import Document

    t0 = time.perf_counter()
    doc = Document(str(filepath))
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    full_text = "\n".join(paragraphs)

    return {
        "name":      filepath.name,
        "path":      str(filepath),
        "type":      "liteparse",
        "parseTime": round(time.perf_counter() - t0, 3),
        "pages": [{
            "pageNum":   1,
            "width":     612,
            "height":    792,
            "text":      full_text,
            "textItems": [],   # no bounding boxes for DOCX
        }],
    }


# ── PPTX parsing with python-pptx ────────────────────────────────────

def _parse_pptx(filepath: Path) -> dict:
    from pptx import Presentation

    t0 = time.perf_counter()
    prs = Presentation(str(filepath))
    pages_data = []

    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs).strip()
                    if line:
                        texts.append(line)
        pages_data.append({
            "pageNum":   i + 1,
            "width":     612,
            "height":    792,
            "text":      "\n".join(texts),
            "textItems": [],   # no bounding boxes for PPTX
        })

    return {
        "name":      filepath.name,
        "path":      str(filepath),
        "type":      "liteparse",
        "parseTime": round(time.perf_counter() - t0, 3),
        "pages":     pages_data,
    }


def _parse_single_file(filepath: Path, _dpi: int) -> dict:
    ext = filepath.suffix.lower()
    if ext in PDF_EXT:
        return _parse_pdf(filepath)
    if ext in DOCX_EXT:
        return _parse_docx(filepath)
    if ext in PPTX_EXT:
        return _parse_pptx(filepath)
    raise ValueError(f"Unsupported format: {ext}")


# ── Parse-Only Mode ──────────────────────────────────────────────────

def run_parse_only(args):
    data_dir = Path(args.dir)
    pdf_files, docx_files, pptx_files, plaintext_files = discover_files(data_dir, args.max_files)
    binary_files = pdf_files + docx_files + pptx_files

    if not binary_files and not plaintext_files:
        print(f"Error: No supported files found in {data_dir}", file=sys.stderr)
        sys.exit(1)

    files_data = []
    total_pages = 0

    if binary_files:
        max_workers = min(args.max_workers, len(binary_files))
        print(f"Parsing {len(binary_files)} files "
              f"({len(pdf_files)} PDF, {len(docx_files)} DOCX, {len(pptx_files)} PPTX) "
              f"with {max_workers} workers...", file=sys.stderr)
        t0_all = time.perf_counter()

        with ThreadPoolExecutor(max_workers=max_workers) as executor:
            futures = {
                executor.submit(_parse_single_file, fp, args.dpi): fp
                for fp in binary_files
            }
            for future in as_completed(futures):
                fp = futures[future]
                try:
                    fd = future.result()
                    total_pages += len(fd["pages"])
                    files_data.append(fd)
                    print(f"  ✓ {fp.name}: {fd['parseTime']}s "
                          f"({len(fd['pages'])} pages)", file=sys.stderr)
                except Exception as e:
                    print(f"  ✗ {fp.name}: FAILED ({e})", file=sys.stderr)

        print(f"Parsed in {time.perf_counter() - t0_all:.2f}s total", file=sys.stderr)

    for fp in plaintext_files:
        print(f"Reading {fp.name}...", file=sys.stderr, end=" ", flush=True)
        try:
            text = fp.read_text(errors="replace")
        except Exception as e:
            print(f"FAILED ({e})", file=sys.stderr)
            continue
        files_data.append({"name": fp.name, "path": str(fp), "type": "plaintext", "text": text})
        print("done", file=sys.stderr)

    output_data = {
        "data_dir": str(data_dir),
        "files": files_data,
        "summary": {
            "total_files":     len(files_data),
            "total_pages":     total_pages,
            "pdf_files":       len(pdf_files),
            "docx_files":      len(docx_files),
            "pptx_files":      len(pptx_files),
            "plaintext_files": len(plaintext_files),
        },
    }

    if args.output == "-":
        print(json.dumps(output_data, ensure_ascii=False))
    else:
        out = Path(args.output)
        out.parent.mkdir(parents=True, exist_ok=True)
        out.write_text(json.dumps(output_data, ensure_ascii=False))
        print(f"\nParsed JSON → {out}", file=sys.stderr)
        print(f"Summary: {output_data['summary']}", file=sys.stderr)


# ── PDF page renderer (PyMuPDF) ───────────────────────────────────────

def _render_pdf_page(filepath: Path, page_num: int, dpi: int) -> dict | None:
    """Screenshot + word boxes for a single PDF page. Returns None on failure."""
    import fitz

    try:
        pdf = fitz.open(str(filepath))
        if page_num < 1 or page_num > len(pdf):
            print(f"  Page {page_num} out of range for {filepath.name}", file=sys.stderr)
            pdf.close()
            return None

        page = pdf[page_num - 1]

        mat = fitz.Matrix(dpi / 72, dpi / 72)
        pix = page.get_pixmap(matrix=mat)
        img_b64 = base64.b64encode(pix.tobytes("png")).decode()

        words = page.get_text("words")
        text_items = [
            {"text": w[4], "x": w[0], "y": w[1], "width": w[2]-w[0], "height": w[3]-w[1]}
            for w in words
        ]

        result = {
            "pageWidth":  page.rect.width,
            "pageHeight": page.rect.height,
            "textItems":  text_items,
            "image":      img_b64,
            "pageText":   page.get_text("text"),
        }
        pdf.close()
        return result
    except Exception as e:
        print(f"  ✗ render failed for {filepath.name} p{page_num}: {e}", file=sys.stderr)
        return None


# ── Markdown to HTML ──────────────────────────────────────────────────

def markdown_to_html(text):
    lines = text.split("\n")
    html_parts = []
    in_list = None
    in_paragraph = False

    for line in lines:
        stripped = line.strip()

        if in_list and not re.match(r"^[-*]\s", stripped) and not re.match(r"^\d+\.\s", stripped):
            html_parts.append(f"</{in_list}>")
            in_list = None

        if not stripped:
            if in_paragraph:
                html_parts.append("</p>")
                in_paragraph = False
            continue

        m = re.match(r"^(#{1,6})\s+(.+)$", stripped)
        if m:
            if in_paragraph:
                html_parts.append("</p>")
                in_paragraph = False
            level = len(m.group(1))
            content = _inline_md(html_module.escape(m.group(2)))
            html_parts.append(f"<h{level}>{content}</h{level}>")
            continue

        m = re.match(r"^[-*]\s+(.+)$", stripped)
        if m:
            if in_paragraph:
                html_parts.append("</p>")
                in_paragraph = False
            if in_list != "ul":
                if in_list:
                    html_parts.append(f"</{in_list}>")
                html_parts.append("<ul>")
                in_list = "ul"
            html_parts.append(f"<li>{_inline_md(html_module.escape(m.group(1)))}</li>")
            continue

        m = re.match(r"^\d+\.\s+(.+)$", stripped)
        if m:
            if in_paragraph:
                html_parts.append("</p>")
                in_paragraph = False
            if in_list != "ol":
                if in_list:
                    html_parts.append(f"</{in_list}>")
                html_parts.append("<ol>")
                in_list = "ol"
            html_parts.append(f"<li>{_inline_md(html_module.escape(m.group(1)))}</li>")
            continue

        content = _inline_md(html_module.escape(stripped))
        if not in_paragraph:
            html_parts.append("<p>")
            in_paragraph = True
        else:
            html_parts.append(" ")
        html_parts.append(content)

    if in_paragraph:
        html_parts.append("</p>")
    if in_list:
        html_parts.append(f"</{in_list}>")

    return "\n".join(html_parts)


def _inline_md(text):
    text = re.sub(r"\*\*(.+?)\*\*", r"<strong>\1</strong>", text)
    text = re.sub(r"\*(.+?)\*", r"<em>\1</em>", text)
    text = re.sub(r"`(.+?)`", r"<code>\1</code>", text)
    text = re.sub(
        r"\[(\d+)\]",
        r'<a class="cite-ref" href="#citation-\1" data-cite="\1">[\1]</a>',
        text,
    )
    return text


# ── Generate Mode ─────────────────────────────────────────────────────

def run_generate(args):
    answer_data = json.loads(Path(args.answer_file).read_text())
    question             = answer_data["question"]
    answer_text          = answer_data["answer"]
    citations            = answer_data["citations"]
    jurisdictions_covered = answer_data.get("jurisdictions_covered", [])

    data_dir = Path(args.dir)
    dpi      = args.dpi

    # Render cited pages — deduplicate by (file, page)
    page_cache: dict = {}
    print(f"Rendering {len(citations)} cited pages...", file=sys.stderr)

    for cit in citations:
        filename = cit["file"]
        page_num = cit.get("page", 0)
        key = (filename, page_num)
        if key in page_cache or page_num == 0:
            continue

        filepath = data_dir / filename
        if not filepath.exists():
            print(f"  Warning: {filepath} not found", file=sys.stderr)
            continue

        ext = filepath.suffix.lower()
        print(f"  {filename} p{page_num}...", file=sys.stderr, end=" ", flush=True)

        if ext in PDF_EXT:
            rendered = _render_pdf_page(filepath, page_num, dpi)
            if rendered:
                page_cache[key] = rendered
                print("done", file=sys.stderr)
            else:
                print("failed", file=sys.stderr)
        else:
            # Non-PDF: text-only citation (no screenshot)
            print("(text-only — not a PDF)", file=sys.stderr)

    # Build citation data list
    citation_data = []
    for cit in citations:
        filename       = cit["file"]
        page_num       = cit.get("page", 0)
        jurisdiction   = cit.get("jurisdiction", "")
        reg_body       = cit.get("regulatory_body", "")
        key            = (filename, page_num)

        base = {
            "file":             filename,
            "page":             page_num,
            "quote":            cit["quote"],
            "relevance":        cit.get("relevance", ""),
            "jurisdiction":     jurisdiction,
            "jurisdictionName": JURISDICTION_NAMES.get(jurisdiction, jurisdiction),
            "jurisdictionFlag": JURISDICTION_FLAGS.get(jurisdiction, ""),
            "regulatoryBody":   reg_body,
        }

        if page_num == 0:
            file_text = ""
            fp = data_dir / filename
            if fp.exists():
                try:
                    file_text = fp.read_text(errors="replace")
                except Exception:
                    pass
            citation_data.append({**base, "type": "plaintext", "pageText": file_text})

        elif key in page_cache:
            cached = page_cache[key]
            citation_data.append({
                **base,
                "type":      "liteparse",
                "pageWidth":  cached["pageWidth"],
                "pageHeight": cached["pageHeight"],
                "textItems":  cached["textItems"],
                "pageText":   cached["pageText"],
                "image":      cached["image"],
            })
        else:
            citation_data.append({**base, "type": "text-only", "pageText": ""})

    # Build template vars
    output_dir_resolved = Path(args.output).resolve()
    data_dir_resolved   = data_dir.resolve()
    try:
        rel_data = os.path.relpath(data_dir_resolved, output_dir_resolved)
    except ValueError:
        rel_data = str(data_dir_resolved)

    images_html_parts = []
    citations_json    = []
    for i, cit in enumerate(citation_data):
        img_id = None
        if cit.get("image"):
            img_id = f"cite-img-{i}"
            images_html_parts.append(
                f'<img id="{img_id}" class="hidden-img" '
                f'src="data:image/png;base64,{cit["image"]}" />'
            )
        slim = {k: v for k, v in cit.items() if k != "image"}
        slim["imgId"] = img_id
        if cit["type"] == "liteparse" and cit["file"].lower().endswith(".pdf"):
            slim["pdfPath"] = f"{rel_data}/{cit['file']}"
        citations_json.append(slim)

    unique_docs  = set(c["file"] for c in citation_data)
    unique_pages = set((c["file"], c["page"]) for c in citation_data if c["page"] > 0)

    jur_display = [
        {
            "code":  code,
            "name":  JURISDICTION_NAMES.get(code, code),
            "flag":  JURISDICTION_FLAGS.get(code, ""),
        }
        for code in jurisdictions_covered
    ]

    answer_html = markdown_to_html(answer_text)

    template_path = Path(args.skill_dir) / "templates" / "report.html"
    if not template_path.exists():
        print(f"Error: Template not found at {template_path}", file=sys.stderr)
        sys.exit(1)

    timestamp = datetime.now().strftime("%Y-%m-%d %H:%M")
    html_out = template_path.read_text()
    html_out = html_out.replace("{{QUESTION}}",          html_module.escape(question))
    html_out = html_out.replace("{{ANSWER_HTML}}",       answer_html)
    html_out = html_out.replace("{{TIMESTAMP}}",         timestamp)
    html_out = html_out.replace("{{NUM_DOCS}}",          str(len(unique_docs)))
    html_out = html_out.replace("{{NUM_PAGES}}",         str(len(unique_pages)))
    html_out = html_out.replace("{{NUM_CITATIONS}}",     str(len(citations_json)))
    html_out = html_out.replace("{{CITATIONS_JSON}}",    json.dumps(citations_json))
    html_out = html_out.replace("{{IMAGES_HTML}}",       "\n".join(images_html_parts))
    html_out = html_out.replace("{{DPI}}",               str(dpi))
    html_out = html_out.replace("{{JURISDICTIONS_JSON}}", json.dumps(jur_display))

    out_dir  = Path(args.output)
    out_dir.mkdir(parents=True, exist_ok=True)
    fname    = f"sankofa-report-{datetime.now().strftime('%Y-%m-%d-%H%M%S')}.html"
    out_path = out_dir / fname
    out_path.write_text(html_out)

    size_mb = out_path.stat().st_size / 1024 / 1024
    print(f"\n✓ Report → {out_path.resolve()}", file=sys.stderr)
    print(f"  {size_mb:.1f} MB | {len(citations_json)} citations | "
          f"Jurisdictions: {', '.join(jurisdictions_covered) or 'n/a'}", file=sys.stderr)

    webbrowser.open(f"file://{out_path.resolve()}")
    print(str(out_path.resolve()))


# ── Main ──────────────────────────────────────────────────────────────

def main():
    ap = argparse.ArgumentParser(description="Sankofa Report Generator")
    ap.add_argument("--skill-dir",   required=True)
    ap.add_argument("--dir",         required=True)
    ap.add_argument("--parse-only",  action="store_true")
    ap.add_argument("--answer-file")
    ap.add_argument("--output",      default="sankofa_output/")
    ap.add_argument("--dpi",         type=int, default=150)
    ap.add_argument("--max-files",   type=int, default=50)
    ap.add_argument("--max-workers", type=int, default=8)
    args = ap.parse_args()

    if args.parse_only:
        run_parse_only(args)
    elif args.answer_file:
        run_generate(args)
    else:
        print("Error: Specify --parse-only or --answer-file", file=sys.stderr)
        sys.exit(1)


if __name__ == "__main__":
    main()
